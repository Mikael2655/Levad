#!/usr/bin/env python3
"""(Re)génère assets/template.pptx à partir de la présentation Offre_Commerciale.

Remplace les champs dynamiques par des jetons {{…}} (run unique) et applique
les retouches statiques demandées. Le navigateur (export-pptx.js) fait ensuite
la substitution des jetons et clone les lignes de tableau par machine.

Usage : python3 tools/prepare_template.py [SOURCE.pptx] [SORTIE.pptx]
Dépendance : python-pptx
"""
import os
import sys
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
SRC = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, "Offre_Commerciale.pptx")
OUT = sys.argv[2] if len(sys.argv) > 2 else os.path.join(ROOT, "assets", "template.pptx")


def shp(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    raise KeyError(name)


def set_para(paragraph, text):
    """Force un paragraphe (déjà pourvu d'au moins un run) à un seul run = text,
    en gardant le format du 1er run."""
    runs = paragraph.runs
    if not runs:
        return
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def drop_extra_paras(tf, keep):
    for p in tf.paragraphs[keep:]:
        p._p.getparent().remove(p._p)


def drop_para(tf, idx):
    p = tf.paragraphs[idx]
    p._p.getparent().remove(p._p)


def fr_all_text(prs):
    """Force lang="fr-FR" sur tous les runs (et fins de paragraphe) de toutes les
    diapositives : le modèle d'origine est tagué en anglais (lang="en-US"), ce qui
    fait souligner en rouge la quasi-totalité du texte français par le correcteur
    orthographique de PowerPoint."""
    for s in prs.slides:
        for el in s._element.iter():
            if etree.QName(el).localname in ("rPr", "defRPr", "endParaRPr"):
                el.set("lang", "fr-FR")


def main():
    prs = Presentation(SRC)
    S = prs.slides

    # ---- Slide 1 (page de garde) ----
    s1 = S[0]
    set_para(shp(s1, "TextBox 13").text_frame.paragraphs[0], "{{DATE}}")
    tb14 = shp(s1, "TextBox 14").text_frame
    set_para(tb14.paragraphs[0], "{{REP_NAME}}")
    set_para(tb14.paragraphs[1], "{{REP_TITLE}}")
    set_para(tb14.paragraphs[2], "{{REP_PHONELINE}}")
    set_para(tb14.paragraphs[3], "{{REP_EMAIL}}")
    drop_para(tb14, 4)   # 5e paragraphe (vide, sans alignement) : inutilisé
    # bas à droite : retire la ligne téléphone (pas juste vidée -> plus d'espace vide),
    # site = www.levad.fr (statique), aligné à la même hauteur que le bloc commercial
    tb15shape = shp(s1, "TextBox 15")
    tb15 = tb15shape.text_frame
    drop_para(tb15, 2)
    set_para(tb15.paragraphs[2], "www.levad.fr")
    tb15shape.top = shp(s1, "TextBox 14").top

    # ---- Slide 4 (lettre) ----
    s4 = S[3]
    tb5 = shp(s4, "TextBox 5").text_frame
    p0 = tb5.paragraphs[0]
    p0.runs[0].text = "Paris, le "
    p0.runs[1].text = "{{DATE_LONG}}"
    for r in p0.runs[2:]:
        r._r.getparent().remove(r._r)
    set_para(tb5.paragraphs[1], "{{CLIENT_CONTACT}}")
    set_para(shp(s4, "TextBox 6").text_frame.paragraphs[0], "{{MACHINE_HEADLINE}}")
    tb7 = shp(s4, "TextBox 7").text_frame
    set_para(tb7.paragraphs[0], "{{REP_NAME}}")
    set_para(tb7.paragraphs[1], "{{REP_TITLE}}")
    set_para(tb7.paragraphs[2], "{{REP_PHONELINE}}")
    set_para(tb7.paragraphs[3], "")           # ancien 2e téléphone : supprimé
    set_para(tb7.paragraphs[4], "{{REP_EMAIL}}")
    tb9 = shp(s4, "TextBox 9").text_frame
    set_para(tb9.paragraphs[0], "{{CLIENT_NAME}}")
    set_para(tb9.paragraphs[1], "{{CLIENT_ADDR1}}")
    set_para(tb9.paragraphs[2], "{{CLIENT_ADDR2}}")

    # ---- Slide 26 (tableaux SA / SP) ----
    s26 = S[25]
    set_para(shp(s26, "TextBox 5").text_frame.paragraphs[0], "SITUATION ACTUELLE € HT / {{PER_UNIT}}")
    set_para(shp(s26, "TextBox 6").text_frame.paragraphs[0], "SOLUTION PROPOSEE € HT / {{PER_UNIT}}")
    tables = [s.table for s in s26.shapes if s.has_table]
    # largeurs de colonnes : Type Machine et Loyer Trimestriel élargies (référence
    # sur 1 ligne), colonnes peu utiles (forfaits, toujours "0") réduites — largeur
    # totale du tableau quasi inchangée.
    col_widths = [2450000, 1300000, 2090000, 1300000, 1300000, 1600000,
                  900000, 900000, 1300000, 1300000, 1443087, 1443087]
    for tbl in tables:
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = Emu(w)
        # "Service Pass" -> "Abonnements et services" (total abonnements + services)
        set_para(tbl.rows[1].cells[5].text_frame.paragraphs[0], "Abonnements et services")
    for tbl, pfx in zip(tables, ["SA", "SP"]):
        cells = tbl.rows[3].cells
        mapping = {0: "TYPE", 1: "FIN", 2: "LOYER", 3: "VNB", 4: "VCOUL",
                   5: "PASS", 6: "FNB", 7: "FCOUL", 8: "CCNB", 9: "CCCOUL", 10: "CE"}
        for ci, key in mapping.items():
            tf = cells[ci].text_frame
            set_para(tf.paragraphs[0], "{{%s_%s}}" % (pfx, key))
            drop_extra_paras(tf, 1)
        c11 = cells[11].text_frame           # total : plus de "soit …/mois"
        set_para(c11.paragraphs[0], "{{%s_TOTAL}}" % pfx)
        drop_extra_paras(c11, 1)

    # ---- Slide 27 (conditions financières) ----
    s27 = S[26]
    tb8 = shp(s27, "TextBox 8").text_frame
    set_para(tb8.paragraphs[1], " Durée : {{DUR_TRIM}} trimestres")
    set_para(tb8.paragraphs[2], " Périodicité {{PER_ADJ}}, terme à échoir")
    # référence machine : une seule ligne (plus de 2e ligne vide), centrée verticalement
    tb17shape = shp(s27, "TextBox 17")
    tb17 = tb17shape.text_frame
    drop_para(tb17, 2)
    drop_para(tb17, 1)
    set_para(tb17.paragraphs[0], "{{PROP_MACHINE_1}}")
    try:
        tb17.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    # loyer : une seule ligne (montant seul, plus de libellé ni de 2e loyer vide),
    # centrée verticalement
    tb26shape = shp(s27, "TextBox 26")
    tb26 = tb26shape.text_frame
    drop_para(tb26, 2)
    drop_para(tb26, 1)
    set_para(tb26.paragraphs[0], "{{PROP_LOYER_1}} € HT")
    try:
        tb26.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    # ---- Slide 28 (contrat de service maintenance) ----
    s28 = S[27]
    tb10 = shp(s28, "TextBox 10").text_frame
    set_para(tb10.paragraphs[0], " Coût Copie N&B {{PROP_MACHINE_1}} : {{SUM_CC_NB}} € HT")
    set_para(tb10.paragraphs[1], " Coût Copie Couleur {{PROP_MACHINE_1}} : {{SUM_CC_COUL}} € HT")

    # ---- Slide 30 (e-maintenance) ----
    s30 = S[29]
    set_para(shp(s30, "TextBox 10").text_frame.paragraphs[0], "Service E-Maintenance : {{EMAINT_VAL}}")

    # ---- Slide 31 (bon pour accord) ----
    s31 = S[30]
    tb10 = shp(s31, "TextBox 10").text_frame
    set_para(tb10.paragraphs[0], "{{REP_NAME}}")
    set_para(tb10.paragraphs[1], "{{REP_TITLE}}")
    set_para(tb10.paragraphs[2], "{{REP_PHONELINE}}")
    set_para(tb10.paragraphs[3], "{{REP_EMAIL}}")
    tb11 = shp(s31, "TextBox 11").text_frame
    set_para(tb11.paragraphs[0], " Valeur {{PER_ADJ_CAP}} : {{SUM_VALEUR}} € HT")
    set_para(tb11.paragraphs[1], " Durée du leasing : {{DUR_TRIM}} Trimestres")
    set_para(tb11.paragraphs[3], " Coût à la page en Noir et Blanc : {{SUM_CC_NB}} € HT")
    set_para(tb11.paragraphs[4], " Coût à la page en Couleur : {{SUM_CC_COUL}} € HT")
    # cadre "date, cachet, signature" vide (Picture 4, simple rectangle) : supprimé —
    # seul le cadre "bon pour accord" (TextBox 12) reste, agrandi, aligné à la même
    # hauteur que le bloc commercial (TextBox 10), texte plaqué en haut du cadre
    # pour laisser la place à la signature du client en dessous.
    pic4 = shp(s31, "Picture 4")
    pic4._element.getparent().remove(pic4._element)
    box = shp(s31, "TextBox 12")
    rep = shp(s31, "TextBox 10")
    box.left = Inches(14.6); box.top = rep.top
    box.width = Inches(5.0); box.height = Inches(2.0)
    box.text_frame.word_wrap = True
    try:
        box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    except Exception:
        pass
    ln = box.line
    ln.color.rgb = RGBColor(0x2F, 0x52, 0x33)
    ln.width = Pt(1.5)

    # ---- Réordonne les diapositives : la 27 (conditions financières avec la
    # machine proposée) prend la place de la 25 (texte générique), qui est
    # supprimée. ----
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    a, b = ids[24], ids[26]     # anciennes diapositives 25 et 27 (0-indexées)
    sldIdLst.remove(a)
    sldIdLst.remove(b)
    sldIdLst.insert(24, b)      # l'ancienne 27 prend la position 25
    sldIdLst.insert(26, a)      # l'ancienne 25 prend la position 27
    cur = list(sldIdLst)
    to_delete = cur[26]         # position 27 = ancienne diapo 25 : supprimée
    rId = to_delete.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(to_delete)

    # Correcteur orthographique : tout le modèle est tagué en anglais (en-US),
    # ce qui souligne en rouge la quasi-totalité du texte français.
    fr_all_text(prs)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("Gabarit écrit :", OUT)
    Presentation(OUT)  # vérifie la réouverture


if __name__ == "__main__":
    main()
